"""多格式文档加载服务

按文件扩展名读取文档内容并返回纯文本：
- txt / md: 直接按 UTF-8 读取
- pdf: 使用 pypdf 提取文本
- pptx: 使用 python-pptx 提取文本框和表格内容
- docx: 使用 python-docx 提取段落和表格内容
- 图片(png/jpg/jpeg/bmp/webp): 走多模态模型(qwen-vl-max)生成文字描述
"""

import base64
from pathlib import Path

from langchain_core.messages import HumanMessage
from loguru import logger

from app.config import config
from app.core.llm_factory import llm_factory

# 支持的文件扩展名（小写，不含点），上传白名单与目录索引共用
SUPPORTED_EXTENSIONS = [
    "txt",
    "md",
    "pdf",
    "pptx",
    "docx",
    "png",
    "jpg",
    "jpeg",
    "bmp",
    "webp",
]

# 图片扩展名 -> MIME 类型
IMAGE_EXTENSIONS = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "bmp": "image/bmp",
    "webp": "image/webp",
}

# 图片识别提示词：要求输出结构化描述，便于后续向量检索
_IMAGE_DESCRIBE_PROMPT = (
    "你是一个文档处理助手。请详细描述这张图片的内容，包括："
    "图片里的全部文字、表格数据、图表信息、操作截图中的步骤和按钮名称等，"
    "输出为纯文本，便于后续做检索。不要发表评论，直接输出内容描述。"
)


class DocumentLoaderService:
    """文档加载服务 - 按扩展名分发解析，输出纯文本内容"""

    def load_document(self, file_path: str) -> str:
        """
        读取文档内容（返回文本）

        Args:
            file_path: 文件路径

        Returns:
            str: 文档文本内容

        Raises:
            ValueError: 文件不存在或扩展名不支持
        """
        path = Path(file_path)
        if not path.exists() or not path.is_file():
            raise ValueError(f"文件不存在: {file_path}")

        ext = path.suffix.lower().lstrip(".")

        if ext in ("txt", "md"):
            content = path.read_text(encoding="utf-8")
            logger.info(f"读取文本文件: {path}, 内容长度: {len(content)} 字符")
            return content
        if ext == "pdf":
            return self._load_pdf(path)
        if ext == "pptx":
            return self._load_pptx(path)
        if ext == "docx":
            return self._load_docx(path)
        if ext in IMAGE_EXTENSIONS:
            return self._describe_image(path)

        raise ValueError(f"不支持的文件格式: {file_path}，支持: {', '.join(SUPPORTED_EXTENSIONS)}")

    def _load_pdf(self, path: Path) -> str:
        """使用 pypdf 提取 PDF 文本"""
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            parts = []
            for i, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                text = text.strip()
                if text:
                    parts.append(f"--- Page {i} ---\n{text}")
            content = "\n\n".join(parts)
            logger.info(f"PDF 解析完成: {path}, {len(reader.pages)} 页, {len(content)} 字符")
            return content
        except ImportError:
            raise RuntimeError("缺少 pypdf 依赖，请执行: uv add pypdf") from None

    def _load_pptx(self, path: Path) -> str:
        """使用 python-pptx 提取幻灯片文本和表格"""
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    texts.extend(self._extract_shape_text(shape))
                # 去重并保留顺序（同一文本框可能被遍历多次）
                seen = set()
                slide_texts = []
                for t in texts:
                    if t not in seen:
                        seen.add(t)
                        slide_texts.append(t)
                if slide_texts:
                    parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))
            content = "\n\n".join(parts)
            logger.info(f"PPTX 解析完成: {path}, {len(prs.slides)} 页, {len(content)} 字符")
            return content
        except ImportError:
            raise RuntimeError("缺少 python-pptx 依赖，请执行: uv add python-pptx") from None

    def _extract_shape_text(self, shape) -> list:
        """递归提取 shape 中的文本（文本框、表格、组合图形）"""
        texts = []
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip(" |"):
                    texts.append(row_text)
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                # 组合图形：递归提取子图形
                for sub in shape.shapes:
                    texts.extend(self._extract_shape_text(sub))
        except ImportError:
            pass
        return texts

    def _load_docx(self, path: Path) -> str:
        """使用 python-docx 提取段落和表格文本"""
        try:
            from docx import Document

            doc = Document(str(path))
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text.strip())
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        parts.append(row_text)
            content = "\n".join(parts)
            logger.info(f"DOCX 解析完成: {path}, {len(content)} 字符")
            return content
        except ImportError:
            raise RuntimeError("缺少 python-docx 依赖，请执行: uv add python-docx") from None

    def _describe_image(self, path: Path) -> str:
        """调用多模态模型生成图片文字描述"""
        ext = path.suffix.lower().lstrip(".")
        mime = IMAGE_EXTENSIONS[ext]

        image_data = base64.b64encode(path.read_bytes()).decode("utf-8")
        data_url = f"data:{mime};base64,{image_data}"

        # 复用 OpenAI 兼容客户端，只换模型为视觉模型
        llm = llm_factory.create_chat_model(
            model=config.vision_model,
            temperature=0.1,
            streaming=False,
        )
        msg = HumanMessage(
            content=[
                {"type": "text", "text": _IMAGE_DESCRIBE_PROMPT},
                {"type": "image_url", "image_url": {"url": data_url}},
            ]
        )
        logger.info(f"开始识别图片: {path}, 大小: {len(image_data) // 1024} KB, 模型: {config.vision_model}")
        resp = llm.invoke([msg])
        content = (resp.content or "").strip()
        if not content:
            raise RuntimeError(f"图片识别返回空内容: {path}")
        logger.info(f"图片识别完成: {path}, 描述长度: {len(content)} 字符")
        return content


# 全局单例
document_loader_service = DocumentLoaderService()
